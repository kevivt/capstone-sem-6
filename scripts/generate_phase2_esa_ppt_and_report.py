from __future__ import annotations

import datetime as dt
import json
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt
from pptx import Presentation


PROJECT_DIR = Path(__file__).resolve().parent.parent


def _load_artifacts_summary() -> dict:
    artifacts_dir = PROJECT_DIR / "artifacts"
    features = {}
    for disease in ["ckd", "hypertension", "diabetes"]:
        fp = artifacts_dir / f"{disease}_features.json"
        if fp.exists():
            features[disease] = json.loads(fp.read_text(encoding="utf-8"))
        else:
            features[disease] = []

    risk_profile = {}
    rp = artifacts_dir / "risk_thresholds_and_factors.json"
    if rp.exists():
        risk_profile = json.loads(rp.read_text(encoding="utf-8"))

    return {"features": features, "risk_profile": risk_profile}


def _set_doc_defaults(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2)
    section.right_margin = Cm(2)

    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(12)

    pf = style.paragraph_format
    pf.line_spacing = 1.5
    pf.space_after = Pt(0)


def _add_chapter(doc: Document, number: int, title: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"CHAPTER {number}\n{title.upper()}")
    run.bold = True
    run.font.size = Pt(18)


def _add_section(doc: Document, sec: str, title: str) -> None:
    p = doc.add_paragraph()
    run = p.add_run(f"{sec} {title}")
    run.bold = True
    run.font.size = Pt(16)


def _add_subsection(doc: Document, sec: str, title: str) -> None:
    p = doc.add_paragraph()
    run = p.add_run(f"{sec} {title}")
    run.bold = True
    run.font.size = Pt(14)


def _add_body(doc: Document, text: str) -> None:
    for para in (text or "").split("\n"):
        t = para.strip()
        if not t:
            doc.add_paragraph("")
        else:
            p = doc.add_paragraph(t)
            p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY


def build_report(output_path: Path) -> Path:
    meta = _load_artifacts_summary()
    feats = meta["features"]
    profiles = meta["risk_profile"].get("diseases", {})

    doc = Document()
    _set_doc_defaults(doc)

    # Minimal placeholder front matter (template expects names/SRNs; keep placeholders).
    title = "Disease Risk Prediction and Disease-Aware Diet Recommendation System"

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Dissertation on\n")
    r.bold = True
    r.font.size = Pt(14)
    p.add_run(f"“{title}”\n").bold = True
    doc.add_paragraph("")

    _add_body(
        doc,
        "Submitted in partial fulfillment of the requirements for the award of the degree of\n"
        "Bachelor of Technology in Computer Science & Engineering\n"
        "UE23CS320B – Capstone Project Phase - 2\n"
        "January – May 2026\n\n"
        "Submitted by:\n"
        "<Name 1> | <SRN1>\n"
        "<Name 2> | <SRN2>\n"
        "<Name 3> | <SRN3>\n"
        "<Name 4> | <SRN4>\n\n"
        "Under the guidance of:\n"
        "<Guide Name>, <Designation>\n"
        "PES University\n",
    )

    doc.add_page_break()

    _add_section(doc, "DECLARATION", "")
    _add_body(
        doc,
        f'We hereby declare that the Capstone Project Phase - 2 entitled “{title}” has been carried out by us '
        "and submitted in partial fulfillment of the course requirements for the award of the degree of "
        "Bachelor of Technology in Computer Science and Engineering of PES University, Bengaluru during "
        "the academic semester January – May 2026. The matter embodied in this report has not been submitted "
        "to any other university or institution for the award of any degree.\n",
    )

    _add_section(doc, "ACKNOWLEDGEMENT", "")
    _add_body(
        doc,
        "We would like to express our gratitude to our project guide and the Department of Computer Science "
        "and Engineering, PES University, for continuous guidance and encouragement throughout Phase - 2.\n"
        "We also thank the Capstone coordinators and panel members for their feedback during reviews.\n",
    )

    _add_section(doc, "ABSTRACT", "")
    _add_body(
        doc,
        "Patients with chronic diseases like Diabetes, Chronic Kidney Disease (CKD), and Hypertension need "
        "to adhere to a strict diet for long-term management. Many existing tools are not user-friendly and do "
        "not adapt to daily lifestyle deviations. This project implements a disease risk prediction + diet "
        "recommendation workflow. Disease-specific models predict risk levels; risk bands are mapped to "
        "nutrition constraints (calorie targets, sodium/potassium limits, and macro guidance). The backend "
        "exposes APIs for prediction, plan generation, meal logging, alerts, and a narrative risk report with "
        "food recommendations from a unified nutrition knowledge base. The next step is a raw patient input "
        "mapping layer to accept user-friendly clinical inputs and deterministically transform them into the "
        "exact feature schema expected by the deployed artifacts.\n",
    )

    doc.add_page_break()

    _add_chapter(doc, 1, "Introduction")
    _add_section(doc, "1.1", "Background")
    _add_body(
        doc,
        "Chronic diseases require long-term diet control and continuous monitoring. A practical system must "
        "combine risk estimation with interpretable diet constraints and provide usable interfaces (backend API "
        "and terminal demo) for validation and iteration.\n",
    )
    _add_section(doc, "1.2", "Project Overview")
    _add_body(
        doc,
        "The system includes: (i) disease risk prediction models (CKD, Hypertension, Diabetes), (ii) a backend "
        "API for prediction, plan recommendation, meal logging and alerting, and (iii) a nutrition knowledge base "
        "to recommend foods consistent with disease constraints.\n",
    )

    _add_chapter(doc, 2, "Problem Definition")
    _add_section(doc, "2.1", "Core Problem")
    _add_body(
        doc,
        "Given patient health indicators, predict disease risk probability and convert risk into actionable "
        "diet targets and safe nutrient limits.\n",
    )
    _add_section(doc, "2.2", "Constraints and Challenges")
    _add_body(
        doc,
        "Key constraints include differing datasets per disease, missing values, categorical fields, and the need "
        "to keep the deployed model feature schema stable. A raw-input user experience must map realistic inputs "
        "to the deployed feature vectors reliably.\n",
    )

    _add_chapter(doc, 3, "Data")
    _add_section(doc, "3.1", "Overview")
    _add_body(
        doc,
        "Datasets used include disease datasets for model training and a unified food knowledge base for nutrition "
        "recommendations.\n",
    )
    _add_section(doc, "3.2", "Dataset and Features (Deployed Artifacts)")
    _add_body(
        doc,
        "This section lists the currently deployed model feature schemas as stored in `artifacts/*_features.json`.\n",
    )
    for disease in ["ckd", "hypertension", "diabetes"]:
        _add_subsection(doc, "3.2." + str(["ckd", "hypertension", "diabetes"].index(disease) + 1), disease.upper())
        cols = feats.get(disease, [])
        _add_body(doc, "Feature columns:\n- " + "\n- ".join([str(c) for c in cols]) + "\n")

    _add_section(doc, "3.3", "Preprocessing Summary (Serving Path)")
    _add_body(
        doc,
        "The current serving artifacts do not include persisted preprocessing transformers. Therefore, Phase - 3 "
        "work will add a deterministic raw-to-feature mapping layer aligned to the deployed feature schema, with "
        "explicit assumptions and safe defaults.\n",
    )

    _add_chapter(doc, 4, "High Level System Design / System Architecture")
    _add_body(
        doc,
        "Major components:\n"
        "- Model artifacts: `artifacts/*_model.joblib` and `artifacts/*_features.json`\n"
        "- Model registry: loads model + feature order; performs single-row inference\n"
        "- Backend API: risk prediction, plan recommendation, meal logging, alerts, risk report\n"
        "- Nutrition KB: unified food dataset used for filtering/scoring recommendations\n",
    )

    _add_chapter(doc, 5, "Implementation and Results (Phase - 2)")
    _add_section(doc, "5.1", "Backend Endpoints Implemented")
    _add_body(
        doc,
        "Implemented endpoints:\n"
        "- GET /health\n"
        "- POST /predict-risk\n"
        "- POST /recommend-plan\n"
        "- POST /log-meal\n"
        "- POST /risk-report\n",
    )
    _add_section(doc, "5.2", "Risk Profiles and Thresholds")
    _add_body(
        doc,
        "Thresholds are stored in `artifacts/risk_thresholds_and_factors.json` and used to classify risk levels "
        "(low/moderate/high). The deployed thresholds are:\n",
    )
    for disease, cfg in profiles.items():
        t = cfg.get("thresholds", {})
        _add_subsection(doc, "5.2." + str(list(profiles.keys()).index(disease) + 1), disease)
        _add_body(
            doc,
            f"Model: {cfg.get('model')}\n"
            f"Moderate threshold: {t.get('moderate')}\n"
            f"High threshold: {t.get('high')}\n",
        )

    _add_section(doc, "5.3", "Diet Calibration (Current)")
    _add_body(
        doc,
        "Diet plans are calibrated using disease type and risk band to set calories and sodium/potassium limits. "
        "Food recommendations are filtered/scored against these limits.\n",
    )

    _add_chapter(doc, 6, "Conclusion of Capstone Project Phase - 2")
    _add_body(
        doc,
        "Phase - 2 achieved an end-to-end backend flow from model artifacts to risk prediction, risk banding, "
        "diet plan generation, meal logging, alerting, and risk-report narrative generation. Terminal prediction "
        "and backend smoke validation are working.\n",
    )

    _add_chapter(doc, 7, "Plan of Work for Capstone Project Phase - 3")
    _add_body(
        doc,
        "Planned Phase - 3 work:\n"
        "- Raw patient input schemas per disease (human-readable clinical fields)\n"
        "- Deterministic raw-to-feature mapping aligned to deployed `*_features.json`\n"
        "- Backend endpoint support for raw-input prediction (backward compatible)\n"
        "- Risk-factor surfacing where possible and tighter diet calibration per disease/risk band\n"
        "- Demo samples, tests, and documentation updates\n",
    )

    _add_chapter(doc, 8, "References / Bibliography")
    _add_body(
        doc,
        "[1] UCI Machine Learning Repository, Chronic Kidney Disease dataset.\n"
        "[2] Framingham Heart Study dataset (cardiovascular risk factors).\n"
        "[3] Pima Indians Diabetes dataset.\n"
        "[4] USDA / combined nutrition datasets used in unified food knowledge base.\n",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return output_path


def build_ppt(output_path: Path) -> Path:
    template = PROJECT_DIR / "esa templates" / "Capstone Project Phase 2 ESA_template.pptx"
    prs = Presentation(str(template))

    meta = _load_artifacts_summary()
    profiles = meta["risk_profile"].get("diseases", {})

    title = "Disease Risk Prediction and Disease-Aware Diet Recommendation System"
    project_id = "UE23CS320B"
    guide = "<Guide Name>"
    team = "<Name 1>, <Name 2>, <Name 3>, <Name 4>"

    abstract = (
        "We built an end-to-end system that predicts disease risk for CKD, Hypertension and Diabetes, "
        "and converts risk bands into medically-aware diet targets and food recommendations. "
        "The backend serves prediction, plan generation, meal logging, alerts, and a narrative risk report. "
        "Next: add a raw patient input mapping layer (user-friendly clinical fields → exact model feature schema) "
        "and tighten disease-aware diet calibration."
    )

    def set_slide_text(slide_idx: int, replacements: dict[str, str]) -> None:
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text:
                continue
            txt = shape.text
            for k, v in replacements.items():
                if k in txt:
                    txt = txt.replace(k, v)
            shape.text = txt

    # Slide 1: title
    set_slide_text(
        0,
        {
            "UE23CS320B – Capstone Project Phase-2 ESA": "UE23CS320B – Capstone Project Phase-2 ESA",
            "Project Title   :": f"Project Title   : {title}\nProject ID       : {project_id}\nProject Guide : {guide}\nProject Team  : {team}",
        },
    )

    # Slide 3: abstract
    set_slide_text(
        2,
        {
            "Well defined problem statement.": "Problem: predict disease risk and translate it into safe, actionable diet targets.",
            "Provide a basic introduction of the project": "System: models + backend APIs + nutrition KB + reporting + terminal demo.",
            "Abstract": f"Abstract\n\n{abstract}",
        },
    )

    # Slide 6: analysis/critique
    analysis = (
        "Research gap: most apps provide generic diet advice and do not enforce disease-specific nutrient constraints.\n"
        "Our approach: risk prediction + explicit thresholds + rule-based nutrient limits + food scoring.\n"
        "Real-world feasibility: deployable backend + terminal demo + stored artifacts + reproducible preprocessing pipeline."
    )
    set_slide_text(5, {"Analysis and Critique of Research/Relevance -Application in Real world": "Analysis and Critique of Research / Relevance / Real-world Application", "Analysis and Critique of Research:": analysis})

    # Slide 7: HLD
    hld = (
        "Components:\n"
        "- Artifacts: trained models + feature schemas + thresholds\n"
        "- ModelRegistry: loads artifacts, enforces feature order, predicts risk score\n"
        "- Backend APIs: /predict-risk, /recommend-plan, /log-meal, /risk-report\n"
        "- Nutrition KB: filters + scores foods using sodium/potassium/sugar/phosphorus constraints\n"
        "Data flow: Patient inputs → feature vector → risk score → risk level → diet targets → food list + report"
    )
    set_slide_text(6, {"Provide high-level design view of the system.": hld})

    # Slide 8: methodology
    meth = (
        "Methodology:\n"
        "- Preprocess datasets (cleaning, missing-value handling, encoding, scaling)\n"
        "- Train disease-wise models and save artifacts (joblib + feature lists)\n"
        "- Build risk profiles: thresholds + top factors\n"
        "- Serve predictions via FastAPI; generate diet plans from risk levels; recommend foods via nutrition KB\n"
        "Planned next: raw-input schemas + deterministic mapping layer + richer diet calibration."
    )
    set_slide_text(7, {"Technology Selection – Choose programming language, tools, and frameworks.": meth})

    # Slide 9: progress
    progress = (
        "Completed:\n"
        "- Trained and saved disease-risk models and feature schemas\n"
        "- Terminal prediction script working\n"
        "- Backend APIs working + smoke test passing\n"
        "- Risk thresholds/top factors generated\n"
        "- Nutrition KB integration for food recommendations\n\n"
        "In progress / next:\n"
        "- Raw patient-friendly input → model feature mapping layer\n"
        "- Connect risk level + major factors into stricter disease-aware diet calibration"
    )
    set_slide_text(8, {"What is the project progress so far?": progress})

    # Slide 10: contribution
    contrib = (
        "Individual contribution (draft):\n"
        "- Data preprocessing pipeline + artifact generation\n"
        "- Backend API integration and smoke testing\n"
        "- Terminal demo flow\n"
        "- Diet plan rules + nutrition KB based recommendations\n"
        "- Reports and comparison outputs generation"
    )
    set_slide_text(9, {"Individual Contribution": f"Individual Contribution\n\n{contrib}"})

    # Slide 13: conclusion
    concl = (
        "Conclusion:\n"
        "- End-to-end risk prediction and diet recommendation workflow implemented.\n"
        "- Deployed artifacts include feature schemas and calibrated thresholds.\n"
        "- Next phase focuses on raw clinical inputs and stronger disease-aware diet calibration."
    )
    set_slide_text(12, {"Summarize the key points.": concl})

    # Slide 14: references
    refs = (
        "[1] UCI ML Repository, Chronic Kidney Disease dataset.\n"
        "[2] Framingham Heart Study dataset.\n"
        "[3] Pima Indians Diabetes dataset.\n"
        "[4] Unified food nutrition knowledge base (merged sources)."
    )
    set_slide_text(13, {"Provide references pertaining to your research according to IEEE format.": refs})

    # Add thresholds summary into slide 11 (Any other info)
    th_lines = []
    for d, cfg in profiles.items():
        t = cfg.get("thresholds", {})
        th_lines.append(f"{d}: moderate={t.get('moderate')}, high={t.get('high')}, model={cfg.get('model')}")
    extra = "Deployed risk thresholds:\n" + "\n".join(th_lines) if th_lines else "Deployed risk thresholds: (see artifacts/risk_thresholds_and_factors.json)"
    set_slide_text(10, {"Provide any other information you wish to add on.": extra})

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main() -> None:
    out_dir = PROJECT_DIR / "outputs" / "esa_phase2"
    out_dir.mkdir(parents=True, exist_ok=True)

    stamp = dt.datetime.now().strftime("%Y%m%d_%H%M")
    report_path = out_dir / f"UE23CS320B_Phase2_Report_{stamp}.docx"
    ppt_path = out_dir / f"UE23CS320B_Phase2_ESA_{stamp}.pptx"

    build_report(report_path)
    build_ppt(ppt_path)

    print(f"Report: {report_path}")
    print(f"PPT: {ppt_path}")


if __name__ == "__main__":
    main()

