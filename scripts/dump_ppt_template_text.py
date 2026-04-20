from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def main() -> None:
    project_dir = Path(__file__).resolve().parent.parent
    template = project_dir / "esa templates" / "Capstone Project Phase 2 ESA_template.pptx"
    if not template.exists():
        raise FileNotFoundError(template)

    prs = Presentation(str(template))
    out_lines: list[str] = []
    out_lines.append(f"slides={len(prs.slides)}")

    for i, slide in enumerate(prs.slides, start=1):
        out_lines.append("")
        out_lines.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = (shape.text or "").strip()
            if txt:
                out_lines.append(txt.replace("\n", " / "))

    out_path = project_dir / "tmp" / "esa_template_extracted" / "phase2_ppt_template_text.txt"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(out_lines), encoding="utf-8")
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()

